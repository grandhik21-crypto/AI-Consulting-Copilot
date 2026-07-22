import json
import uuid
from pathlib import Path

from pydantic import BaseModel
from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.services.gemini_service import get_model
from app.services.dashboard_store import add_slides

router = APIRouter()

SLIDES_DIR = Path(__file__).resolve().parent.parent.parent / "uploads" / "slides"
SLIDES_DIR.mkdir(parents=True, exist_ok=True)

BG_DARK = RGBColor(0x12, 0x12, 0x12)
BG_CARD = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT = RGBColor(0x63, 0x63, 0xF0)  # indigo-500
WHITE = RGBColor(0xF4, 0xF4, 0xF4)
GRAY = RGBColor(0xA1, 0xA1, 0xAA)
AMBER = RGBColor(0xF5, 0xA6, 0x23)


class SlidesGenerateRequest(BaseModel):
    executive_summary: str


class SlideOutline(BaseModel):
    slide_title: str
    bullet_points: list[str]
    chart_to_include: str | None = None
    speaker_notes: str


class SlidesGenerateResponse(BaseModel):
    slide_deck_id: str
    title: str
    slides: list[SlideOutline]
    filename: str


def _build_pptx(title: str, slides: list[SlideOutline], filename: str) -> Path:
    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 16:9
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]  # blank

    def _set_bg(slide, color=BG_DARK):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1)
        return shape

    def _add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    # --- Title slide ---
    title_slide = prs.slides.add_slide(blank_layout)
    _set_bg(title_slide)

    _add_shape(title_slide, Inches(1), Inches(2.5), Inches(8), Inches(0.04), fill_color=ACCENT)

    _add_textbox(title_slide, Inches(1), Inches(1.2), Inches(8), Inches(1.2),
                 title, font_size=36, color=WHITE, bold=True)

    _add_textbox(title_slide, Inches(1), Inches(2.8), Inches(8), Inches(0.6),
                 "AI Consulting Copilot", font_size=16, color=GRAY)

    _add_textbox(title_slide, Inches(1), Inches(3.6), Inches(8), Inches(0.5),
                 "Generated Slide Deck", font_size=12, color=GRAY)

    _add_shape(title_slide, Inches(1), Inches(4.5), Inches(8), Inches(0.04), fill_color=ACCENT)

    # --- Content slides ---
    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide)

        # Slide number bar at top
        _add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), fill_color=ACCENT)

        # Slide number
        _add_textbox(slide, Inches(0.6), Inches(0.3), Inches(1), Inches(0.4),
                     f"{i + 1}", font_size=12, color=GRAY)

        # Title
        _add_textbox(slide, Inches(1), Inches(0.3), Inches(8), Inches(0.5),
                     slide_data["slide_title"], font_size=28, color=WHITE, bold=True)

        # Accent line under title
        _add_shape(slide, Inches(1), Inches(0.9), Inches(3), Inches(0.03), fill_color=ACCENT)

        # Bullet points
        bullets_text = "\n".join(f"•  {b}" for b in slide_data["bullet_points"])
        _add_textbox(slide, Inches(1), Inches(1.3), Inches(7.5), Inches(2.5),
                     bullets_text, font_size=14, color=WHITE)

        # Chart recommendation
        if slide_data.get("chart_to_include"):
            chart_box = _add_shape(slide, Inches(1), Inches(3.9), Inches(7.5), Inches(0.9),
                                   fill_color=BG_CARD, line_color=ACCENT)
            chart_box.line.width = Pt(1)

            _add_textbox(slide, Inches(1.2), Inches(3.95), Inches(1.5), Inches(0.3),
                         "📊  Chart:", font_size=11, color=ACCENT, bold=True)
            _add_textbox(slide, Inches(2.5), Inches(3.95), Inches(5.5), Inches(0.7),
                         slide_data["chart_to_include"], font_size=11, color=AMBER)

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_data.get("speaker_notes", "")

        # Page footer
        _add_shape(slide, Inches(0), Inches(7.2), Inches(10), Inches(0.04), fill_color=ACCENT)
        _add_textbox(slide, Inches(0.6), Inches(6.7), Inches(8), Inches(0.3),
                     "Confidential  •  AI Consulting Copilot", font_size=8, color=GRAY)

    save_path = SLIDES_DIR / filename
    prs.save(str(save_path))
    return save_path


GENERATE_PROMPT = (
    "You are a senior management consultant creating a professional slide deck outline. "
    "Given the executive summary below, produce a structured slide deck outline.\n\n"
    "For each slide provide:\n"
    "- slide_title: concise, professional title\n"
    "- bullet_points: array of 3-5 key bullet points\n"
    "- chart_to_include: description of a chart or data visualization that would support this slide (or null if not applicable)\n"
    "- speaker_notes: 2-3 sentences the presenter should say\n\n"
    "Output ONLY valid JSON with this exact structure (no markdown, no code fences):\n"
    "{\n"
    '  "title": "Deck Title",\n'
    '  "slides": [\n'
    "    {\n"
    '      "slide_title": "...",\n'
    '      "bullet_points": ["...", "..."],\n'
    '      "chart_to_include": "..." or null,\n'
    '      "speaker_notes": "..."\n'
    "    }\n"
    "  ]\n"
    "}\n\n"
    "Executive Summary:\n{text}"
)


@router.post("/slides/generate", response_model=SlidesGenerateResponse)
async def generate_slides(body: SlidesGenerateRequest):
    if not body.executive_summary.strip():
        raise HTTPException(status_code=400, detail="Executive summary is required")

    model = get_model()

    try:
        raw = model.generate_content(GENERATE_PROMPT.format(text=body.executive_summary[:10000])).text.strip()
        raw = raw.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        data = json.loads(raw)

        title = data.get("title", "Consulting Slide Deck")
        slides = data.get("slides", [])

        deck_id = uuid.uuid4().hex[:12]
        filename = f"slide_deck_{deck_id}.pptx"

        _build_pptx(title, slides, filename)

        add_slides(title, len(slides), filename)

        return SlidesGenerateResponse(
            slide_deck_id=deck_id,
            title=title,
            slides=[SlideOutline(**s) for s in slides],
            filename=filename,
        )

    except json.JSONDecodeError:
        raise HTTPException(status_code=502, detail="Failed to parse slide outline from Gemini response")
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Slide generation failed: {str(e)}")


@router.get("/slides/{filename}")
async def download_slides(filename: str):
    file_path = SLIDES_DIR / filename
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(
        path=str(file_path),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )
